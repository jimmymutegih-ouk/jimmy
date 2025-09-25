from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create a new presentation
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Title Only layout
slide = prs.slides.add_slide(slide_layout)

# Add slide title
title = slide.shapes.title
title.text = "Equity Structure After Investment (KES 7.5M)"

# Add investment summary textbox
left = Inches(0.5)
top = Inches(1.5)
width = Inches(4.5)
height = Inches(1.5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame

# Summary points
points = [
    "Pre-money Valuation: KES 10M",
    "Investment Raised: KES 7.5M",
    "Post-money Valuation: KES 17.5M"
]

# Add each point to the textbox
for point in points:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(14)

# Chart data: ownership distribution
chart_data = ChartData()
chart_data.categories = ['Abdallah (CEO)', 'Jim (Co-founder)', 'Investor']
chart_data.add_series('Ownership', (36, 21, 43))

# Chart location and size
x, y, cx, cy = Inches(5.2), Inches(1.2), Inches(4), Inches(3.5)

# Add pie chart
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

# Show legend
chart.has_legend = True
chart.legend.include_in_layout = False

# Save presentation to a file
prs.save("Fix_Tea_Equity_Structure_Pitch_Deck.pptx")
